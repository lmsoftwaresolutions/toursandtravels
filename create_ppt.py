"""
Nathkrupa Tours & Travels - Project Presentation Generator
Creates a professional, animated PPT with all chapters from the index.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Configuration ───────────────────────────────────────────────
SAVE_PATH = r"C:\Users\NIKITA\Desktop\test1\toursandtravels\Nathkrupa_Tours_Travels_Project_v2.pptx"
DIAGRAM_DIR = r"C:\Users\NIKITA\.gemini\antigravity\brain\ca2895fb-7803-4546-8f23-86f3a2816a3a"

# Colors
BG_DARK     = RGBColor(0x0D, 0x0D, 0x0D)
BG_CARD     = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_RED  = RGBColor(0xE5, 0x3E, 0x3E)
ACCENT_GOLD = RGBColor(0xFF, 0xD7, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_GRAY    = RGBColor(0x88, 0x88, 0x88)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ─── Helpers ─────────────────────────────────────────────────────

def add_entrance_animation(slide, shape, delay_ms=0):
    """Add fade-in animation to a shape using raw XML."""
    import lxml.etree as etree
    
    # Define namespaces
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    
    # Get shape id
    sp_id = shape.shape_id
    
    timing = slide._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}timing')
    if timing is None:
        timing = etree.SubElement(slide._element, '{http://schemas.openxmlformats.org/presentationml/2006/main}timing')
    
    tnLst = timing.find('{http://schemas.openxmlformats.org/presentationml/2006/main}tnLst')
    if tnLst is None:
        tnLst = etree.SubElement(timing, '{http://schemas.openxmlformats.org/presentationml/2006/main}tnLst')
    
    par = tnLst.find('{http://schemas.openxmlformats.org/presentationml/2006/main}par')
    if par is None:
        par = etree.SubElement(tnLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}par')
        cTn_root = etree.SubElement(par, '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn', attrib={
            'id': '1', 'dur': 'indefinite', 'restart': 'never', 'nodeType': 'tmRoot'
        })
        childTnLst = etree.SubElement(cTn_root, '{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')
        seq_par = etree.SubElement(childTnLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}seq', attrib={
            'concurrent': '1', 'nextAc': 'seek'
        })
        seq_cTn = etree.SubElement(seq_par, '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn', attrib={
            'id': '2', 'dur': 'indefinite', 'nodeType': 'mainSeq'
        })
        seq_childTnLst = etree.SubElement(seq_cTn, '{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')
    else:
        cTn_root = par.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cTn')
        childTnLst = cTn_root.find('{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')
        seq_par = childTnLst.find('{http://schemas.openxmlformats.org/presentationml/2006/main}seq')
        seq_cTn = seq_par.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cTn')
        seq_childTnLst = seq_cTn.find('{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')
        if seq_childTnLst is None:
            seq_childTnLst = etree.SubElement(seq_cTn, '{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')

    # Count existing to generate unique IDs
    existing = seq_childTnLst.findall('{http://schemas.openxmlformats.org/presentationml/2006/main}par')
    base_id = 3 + len(existing) * 6
    
    # Build animation par element for fade
    anim_par = etree.SubElement(seq_childTnLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}par')
    anim_cTn = etree.SubElement(anim_par, '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn', attrib={
        'id': str(base_id), 'fill': 'hold'
    })
    stCondLst = etree.SubElement(anim_cTn, '{http://schemas.openxmlformats.org/presentationml/2006/main}stCondLst')
    etree.SubElement(stCondLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}cond', attrib={'delay': '0'})
    
    inner_childTnLst = etree.SubElement(anim_cTn, '{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')
    inner_par = etree.SubElement(inner_childTnLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}par')
    inner_cTn = etree.SubElement(inner_par, '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn', attrib={
        'id': str(base_id + 1), 'presetID': '10', 'presetClass': 'entr', 'presetSubtype': '0',
        'fill': 'hold', 'nodeType': 'afterEffect'
    })
    inner_stCondLst = etree.SubElement(inner_cTn, '{http://schemas.openxmlformats.org/presentationml/2006/main}stCondLst')
    etree.SubElement(inner_stCondLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}cond', attrib={
        'delay': str(delay_ms)
    })
    
    effect_childTnLst = etree.SubElement(inner_cTn, '{http://schemas.openxmlformats.org/presentationml/2006/main}childTnLst')
    
    # AnimEffect for fade
    animEffect = etree.SubElement(effect_childTnLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}animEffect', attrib={
        'transition': 'in', 'filter': 'fade'
    })
    ae_cBhvr = etree.SubElement(animEffect, '{http://schemas.openxmlformats.org/presentationml/2006/main}cBhvr')
    ae_cTn = etree.SubElement(ae_cBhvr, '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn', attrib={
        'id': str(base_id + 2), 'dur': '500'
    })
    ae_tgtEl = etree.SubElement(ae_cBhvr, '{http://schemas.openxmlformats.org/presentationml/2006/main}tgtEl')
    ae_spTgt = etree.SubElement(ae_tgtEl, '{http://schemas.openxmlformats.org/presentationml/2006/main}spTgt', attrib={
        'spid': str(sp_id)
    })
    
    # Set element
    setEl = etree.SubElement(effect_childTnLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}set')
    set_cBhvr = etree.SubElement(setEl, '{http://schemas.openxmlformats.org/presentationml/2006/main}cBhvr')
    set_cTn = etree.SubElement(set_cBhvr, '{http://schemas.openxmlformats.org/presentationml/2006/main}cTn', attrib={
        'id': str(base_id + 3), 'dur': '1', 'fill': 'hold'
    })
    set_stCondLst = etree.SubElement(set_cTn, '{http://schemas.openxmlformats.org/presentationml/2006/main}stCondLst')
    etree.SubElement(set_stCondLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}cond', attrib={'delay': '0'})
    set_tgtEl = etree.SubElement(set_cBhvr, '{http://schemas.openxmlformats.org/presentationml/2006/main}tgtEl')
    etree.SubElement(set_tgtEl, '{http://schemas.openxmlformats.org/presentationml/2006/main}spTgt', attrib={
        'spid': str(sp_id)
    })
    set_attrNameLst = etree.SubElement(set_cBhvr, '{http://schemas.openxmlformats.org/presentationml/2006/main}attrNameLst')
    set_attrName = etree.SubElement(set_attrNameLst, '{http://schemas.openxmlformats.org/presentationml/2006/main}attrName')
    set_attrName.text = 'style.visibility'
    set_to = etree.SubElement(setEl, '{http://schemas.openxmlformats.org/presentationml/2006/main}to')
    set_to_val = etree.SubElement(set_to, '{http://schemas.openxmlformats.org/presentationml/2006/main}strVal', attrib={'val': 'visible'})


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(4)
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=16,
                      color=WHITE, font_name="Calibri", bullet=False, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if bullet:
            p.text = f"▸  {line}"
        else:
            p.text = line

        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.space_before = Pt(2)
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_RED, height=Inches(0.05)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_lines, title_color=ACCENT_RED):
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.line.color.rgb = RGBColor(0x33, 0x33, 0x55)
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # Title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.6), Inches(0.45),
                 title, font_size=18, bold=True, color=title_color)

    # Body
    add_multiline_box(slide, left + Inches(0.3), top + Inches(0.6), width - Inches(0.6),
                      height - Inches(0.8), body_lines, font_size=14, color=LIGHT_GRAY, bullet=True)
    return card


def add_slide_number(slide, num):
    add_text_box(slide, W - Inches(1), H - Inches(0.5), Inches(0.8), Inches(0.3),
                 str(num), font_size=11, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, section_num, section_title):
    set_slide_bg(slide, BG_DARK)
    # Big accent rectangle
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), W, Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_RED
    shape.line.fill.background()

    add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 f"CHAPTER {section_num}", font_size=22, bold=False, color=RGBColor(0xFF, 0xCC, 0xCC),
                 alignment=PP_ALIGN.CENTER, font_name="Calibri Light")
    add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(1.2),
                 section_title, font_size=44, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)


def add_content_slide(slide, title, page_num):
    set_slide_bg(slide, BG_DARK)
    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BG_CARD
    bar.line.fill.background()

    # Red accent
    add_accent_line(slide, Inches(0), Inches(1.1), W)

    # Title
    s = add_text_box(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
                 title, font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
    add_entrance_animation(slide, s, 0)
    add_slide_number(slide, page_num)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

# Decorative red stripe
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), H)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

# Horizontal accent
add_accent_line(slide, Inches(2), Inches(3.2), Inches(9))
add_accent_line(slide, Inches(2), Inches(4.9), Inches(9))

# Title text
s1 = add_text_box(slide, Inches(2), Inches(1.5), Inches(9.5), Inches(1.5),
             "NATHKRUPA", font_size=60, bold=True, color=WHITE,
             alignment=PP_ALIGN.LEFT, font_name="Calibri Light")
s2 = add_text_box(slide, Inches(2), Inches(2.4), Inches(9.5), Inches(1),
             "TOURS & TRAVELS", font_size=54, bold=True, color=ACCENT_RED,
             alignment=PP_ALIGN.LEFT)
# Subtitle
s3 = add_text_box(slide, Inches(2), Inches(3.5), Inches(9.5), Inches(1.2),
             "Tour & Travel Management System\nProject Documentation & Analysis",
             font_size=24, bold=False, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT,
             font_name="Calibri Light")
s4 = add_text_box(slide, Inches(2), Inches(5.2), Inches(9.5), Inches(1),
             "A comprehensive web-based solution for managing trips, bookings,\nvehicles, drivers, payments, and invoicing.",
             font_size=16, bold=False, color=MED_GRAY, alignment=PP_ALIGN.LEFT)

for i, s in enumerate([s1, s2, s3, s4]):
    add_entrance_animation(slide, s, i * 300)

add_slide_number(slide, 1)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 2 — TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "TABLE OF CONTENTS", 2)

toc_left = [
    "1.    Introduction",
    "1.1   Project Objectives",
    "1.2   Existing System & Need",
    "1.3   Scope of Work",
    "1.4   Operating Environment",
    "1.5   Technology Used",
    "1.6   Module Specification",
]
toc_right = [
    "2.    Analysis & Design",
    "2.1   Data Flow Diagram",
    "2.2   Entity Relationship Diagram",
    "2.3   Use Case Diagram",
    "2.4   Class Diagram",
    "2.5   Activity Diagram",
    "2.6   Sequence Diagram",
]

s5 = add_multiline_box(slide, Inches(1.5), Inches(1.8), Inches(5), Inches(4.5),
                  toc_left, font_size=20, color=LIGHT_GRAY)
s6 = add_multiline_box(slide, Inches(7), Inches(1.8), Inches(5), Inches(4.5),
                  toc_right, font_size=20, color=LIGHT_GRAY)
add_entrance_animation(slide, s5, 200)
add_entrance_animation(slide, s6, 500)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 3 — CHAPTER 1 HEADER
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 1, "INTRODUCTION")
add_slide_number(slide, 3)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 4 — 1.1 PROJECT OBJECTIVES
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "1.1  Project Objectives", 4)

objectives = [
    "Digitize and automate the complete tour & travel management workflow.",
    "Streamline trip booking, vehicle allocation, and driver assignment processes.",
    "Provide real-time tracking of payments, expenses, and pending balances.",
    "Generate professional invoices, booking receipts, and quotations instantly.",
    "Maintain a centralized database of customers, vehicles, drivers, and vendors.",
    "Offer a comprehensive dashboard with business analytics and summaries.",
    "Enable multi-vehicle trip management with per-vehicle expense tracking.",
    "Reduce manual paperwork and eliminate data redundancy errors.",
]
s = add_multiline_box(slide, Inches(1), Inches(1.6), Inches(11), Inches(5),
                  objectives, font_size=18, color=LIGHT_GRAY, bullet=True)
add_entrance_animation(slide, s, 300)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 5 — 1.2 EXISTING SYSTEM & NEED
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "1.2  Existing System & Need of System", 5)

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(4.8),
         "⚠  Existing System (Manual)", [
             "All records maintained in paper registers and notebooks.",
             "Trip details, customer info tracked manually — prone to errors.",
             "Payment tracking done via handwritten ledgers.",
             "Invoice & receipt generation is slow and inconsistent.",
             "No centralized data — difficult to retrieve past records.",
             "No analytics or business performance insights.",
         ])

add_card(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(4.8),
         "✓  Proposed System (Automated)", [
             "Fully digital trip management with real-time data entry.",
             "Automated vehicle & driver assignment with availability tracking.",
             "Instant payment recording with advance & balance tracking.",
             "One-click professional invoice and receipt generation.",
             "Centralized PostgreSQL database for all business data.",
             "Dashboard with charts, summaries, and key business metrics.",
         ], title_color=RGBColor(0x4C, 0xAF, 0x50))


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 6 — 1.3 SCOPE OF WORK
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "1.3  Scope of Work", 6)

scope = [
    "User authentication & role-based access control (JWT-based).",
    "Complete CRUD operations for Trips, Customers, Vehicles & Drivers.",
    "Multi-vehicle trip support with individual expense tracking per vehicle.",
    "Payment management — advance payments, partial payments & balances.",
    "Automated invoice & booking receipt generation (print-ready layouts).",
    "Quotation creation and management for prospective clients.",
    "Fuel, maintenance & spare parts tracking for vehicle fleet management.",
    "Vendor management for fuel suppliers, mechanics & service providers.",
    "Comprehensive reporting with trip-wise and period-wise summaries.",
    "Dashboard with real-time business insights and key metrics.",
]
s = add_multiline_box(slide, Inches(1), Inches(1.6), Inches(11), Inches(5.2),
                  scope, font_size=17, color=LIGHT_GRAY, bullet=True)
add_entrance_animation(slide, s, 300)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 7 — 1.4 OPERATING ENVIRONMENT
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "1.4  Operating Environment — Hardware & Software", 7)

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5),
         "🖥  Hardware Requirements", [
             "Processor: Intel i3 / AMD Ryzen 3 or above",
             "RAM: 4 GB minimum (8 GB recommended)",
             "Storage: 20 GB free disk space",
             "Network: Stable Internet connection",
             "Display: 1366×768 minimum resolution",
         ])

add_card(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5),
         "💻  Software Requirements", [
             "OS: Windows 10/11, Linux, or macOS",
             "Browser: Chrome, Firefox, or Edge (latest)",
             "Backend: Python 3.12+ with FastAPI",
             "Frontend: Node.js 18+ with React (Vite)",
             "Database: PostgreSQL 15",
             "Container: Docker & Docker Compose",
             "Web Server: Nginx (reverse proxy)",
         ])


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 8 — 1.5 TECHNOLOGY USED
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "1.5  Technology Used", 8)

# Create 3 cards for tech stack
add_card(slide, Inches(0.4), Inches(1.5), Inches(3.9), Inches(5),
         "🔧  Backend", [
             "Python 3.12",
             "FastAPI 0.110",
             "SQLAlchemy 2.0 (ORM)",
             "Alembic (Migrations)",
             "Pydantic (Validation)",
             "JWT Authentication",
             "Uvicorn / Gunicorn",
         ])

add_card(slide, Inches(4.7), Inches(1.5), Inches(3.9), Inches(5),
         "🎨  Frontend", [
             "React 18 (Vite)",
             "JavaScript (ES6+)",
             "React Router v6",
             "Axios (HTTP Client)",
             "CSS3 (Custom Styling)",
             "Responsive Design",
             "Print-Ready Layouts",
         ])

add_card(slide, Inches(9), Inches(1.5), Inches(3.9), Inches(5),
         "🗄  Database & DevOps", [
             "PostgreSQL 15",
             "Docker & Docker Compose",
             "Nginx (Reverse Proxy)",
             "Alpine Linux Containers",
             "Git Version Control",
             "RESTful API Design",
             "Environment Variables",
         ])


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 9-10 — 1.6 MODULE SPECIFICATION
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "1.6  Module Specification (Part 1)", 9)

add_card(slide, Inches(0.3), Inches(1.5), Inches(4), Inches(2.3),
         "🔐  Authentication", [
             "JWT-based login & registration",
             "Password hashing (bcrypt)",
             "Protected API routes",
         ])

add_card(slide, Inches(4.6), Inches(1.5), Inches(4), Inches(2.3),
         "🚌  Trip Management", [
             "Create, edit, view & delete trips",
             "Multi-vehicle support per trip",
             "Route, KM & duration tracking",
         ])

add_card(slide, Inches(8.9), Inches(1.5), Inches(4), Inches(2.3),
         "👥  Customer Management", [
             "Customer CRUD operations",
             "Name, phone, email & address",
             "Auto-link to trip bookings",
         ])

add_card(slide, Inches(0.3), Inches(4.1), Inches(4), Inches(2.3),
         "🚗  Vehicle Management", [
             "Vehicle registry with type & seats",
             "Vehicle notes & history",
             "Fuel & maintenance logs",
         ])

add_card(slide, Inches(4.6), Inches(4.1), Inches(4), Inches(2.3),
         "👤  Driver Management", [
             "Driver profiles & contact info",
             "Driver salary tracking",
             "Driver expense recording",
         ])

add_card(slide, Inches(8.9), Inches(4.1), Inches(4), Inches(2.3),
         "💰  Payment Management", [
             "Advance & partial payments",
             "Payment mode tracking",
             "Balance auto-calculation",
         ])


# Module Spec Part 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "1.6  Module Specification (Part 2)", 10)

add_card(slide, Inches(0.3), Inches(1.5), Inches(4), Inches(2.3),
         "📄  Invoice & Receipt", [
             "Auto-generated invoices",
             "Booking receipts (print-ready)",
             "Professional PDF-style layouts",
         ])

add_card(slide, Inches(4.6), Inches(1.5), Inches(4), Inches(2.3),
         "📝  Quotation Module", [
             "Create quotations for clients",
             "Pricing items & charge items",
             "Convert to trip bookings",
         ])

add_card(slide, Inches(8.9), Inches(1.5), Inches(4), Inches(2.3),
         "⛽  Fuel & Maintenance", [
             "Fuel purchase tracking",
             "Maintenance scheduling",
             "Spare parts inventory",
         ])

add_card(slide, Inches(0.3), Inches(4.1), Inches(4), Inches(2.3),
         "🏪  Vendor Management", [
             "Vendor profiles & categories",
             "Vendor payment tracking",
             "Fuel supplier management",
         ])

add_card(slide, Inches(4.6), Inches(4.1), Inches(4), Inches(2.3),
         "📊  Reports & Dashboard", [
             "Trip-wise & period-wise reports",
             "Revenue & expense analytics",
             "Dashboard with key metrics",
         ])

add_card(slide, Inches(8.9), Inches(4.1), Inches(4), Inches(2.3),
         "📌  Notes Module", [
             "Dashboard sticky notes",
             "Quick reminders & to-dos",
             "Pin important information",
         ])


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 11 — CHAPTER 2 HEADER
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 2, "ANALYSIS & DESIGN")
add_slide_number(slide, 11)


# ═══════════════════════════════════════════════════════════════════
#  DIAGRAM SLIDES — Helper
# ═══════════════════════════════════════════════════════════════════
def add_diagram_slide(title, image_path, page_num, description=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, title, page_num)

    if description:
        add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5),
                     description, font_size=14, color=MED_GRAY)

    if os.path.exists(image_path):
        # Add image centered
        try:
            from PIL import Image
            img = Image.open(image_path)
            img_w, img_h = img.size
            aspect = img_w / img_h

            max_w = Inches(10)
            max_h = Inches(5)

            if aspect > (max_w / max_h):
                w = max_w
                h = int(w / aspect)
            else:
                h = max_h
                w = int(h * aspect)

            left = int((W - w) / 2)
            top = Inches(1.8) if description else Inches(1.5)

            pic = slide.shapes.add_picture(image_path, left, top, w, h)
            add_entrance_animation(slide, pic, 300)
        except Exception as e:
            add_text_box(slide, Inches(2), Inches(3), Inches(9), Inches(1),
                         f"[Image loading error: {str(e)}]", font_size=16, color=ACCENT_RED)
    else:
        add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(1),
                     f"[Diagram image not found at: {image_path}]",
                     font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    return slide


# ─── Diagram image paths (user-uploaded originals + generated) ───
activity_img = os.path.join(DIAGRAM_DIR, "activity_diagram_1774548569082.png")
sequence_img = os.path.join(DIAGRAM_DIR, "sequence_diagram_1774548930265.png")
context_dfd_img = os.path.join(DIAGRAM_DIR, "media__1774587277547.png")     # Context Level DFD
first_level_dfd_img = os.path.join(DIAGRAM_DIR, "media__1774586714032.png")  # First Level DFD
er_diagram_img = os.path.join(DIAGRAM_DIR, "media__1774586649215.png")       # ER Diagram
use_case_img = os.path.join(DIAGRAM_DIR, "media__1774586740172.png")         # Use Case Diagram
class_diagram_img = os.path.join(DIAGRAM_DIR, "media__1774586757843.png")    # Class Diagram


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 12 — 2.1 DATA FLOW DIAGRAM (Context Level)
# ═══════════════════════════════════════════════════════════════════
add_diagram_slide("2.1  Data Flow Diagram — Context Level", context_dfd_img, 12,
                  "Shows the system as a single process with Admin/User and Database entities.")


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 13 — 2.1 DFD - First Level
# ═══════════════════════════════════════════════════════════════════
add_diagram_slide("2.1  Data Flow Diagram — First Level DFD", first_level_dfd_img, 13,
                  "Decomposes the system into 7 core processes: Login, Manage Trips, Customers, Vehicles, Payments, Invoices, Reports.")


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 14 — 2.2 ER DIAGRAM
# ═══════════════════════════════════════════════════════════════════
add_diagram_slide("2.2  Entity Relationship Diagram", er_diagram_img, 14,
                  "Chen notation ER diagram showing all entities, attributes, relationships and cardinalities.")


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 15 — 2.3 USE CASE DIAGRAM
# ═══════════════════════════════════════════════════════════════════
add_diagram_slide("2.3  Use Case Diagram", use_case_img, 15,
                  "UML Use Case diagram with User (Limited Access) and Admin (Full Access) actors.")


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 16 — 2.4 CLASS DIAGRAM
# ═══════════════════════════════════════════════════════════════════
add_diagram_slide("2.4  Class Diagram", class_diagram_img, 16,
                  "UML Class diagram with attributes, methods, and relationships for all core classes.")


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 17 — 2.5 ACTIVITY DIAGRAM
# ═══════════════════════════════════════════════════════════════════
if activity_img:
    add_diagram_slide("2.5  Activity Diagram", activity_img, 17,
                      "Shows the workflow from Admin Login → Trip Management → Invoice Generation.")
else:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "2.5  Activity Diagram", 17)
    act_lines = [
        "Flow: Start → Admin Login → [Credentials Valid?]",
        "   [No] → Loop back to Login",
        "   [Yes] → Admin Dashboard",
        "   → Manage Trips (Create/Edit)",
        "   → Assign Vehicles & Drivers",
        "   → Manage Customer Bookings",
        "   → Generate Booking Receipt",
        "   → Record Payments & Expenses",
        "   → Generate Invoice / Reports",
        "   → Admin Logout → End",
    ]
    s = add_multiline_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.2),
                      act_lines, font_size=18, color=LIGHT_GRAY, bullet=True)
    add_entrance_animation(slide, s, 300)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 18 — 2.6 SEQUENCE DIAGRAM
# ═══════════════════════════════════════════════════════════════════
if sequence_img:
    add_diagram_slide("2.6  Sequence Diagram", sequence_img, 18,
                      "Shows object interactions: Admin ↔ System ↔ Trip Manager ↔ Customer Manager ↔ Vehicle Manager ↔ Payment System.")
else:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_content_slide(slide, "2.6  Sequence Diagram", 18)
    seq_lines = [
        "Objects: Admin/User, Tours & Travels System, Trip Manager,",
        "              Customer Manager, Vehicle Manager, Payment System",
        "",
        "Interaction Sequence:",
        "1. Admin → System: Login() → Verify Credentials → Login Successful",
        "2. System → Admin: Display Dashboard()",
        "3. Admin → Trip Manager: Create Trip()",
        "4. Trip Manager → Vehicle Manager: Assign Vehicle()",
        "5. Vehicle Manager → Trip Manager: Vehicle Assigned()",
        "6. Admin → Customer Manager: Add Booking()",
        "7. Customer Manager → Admin: Booking Confirmed()",
        "8. Admin → Payment System: Record Payment()",
        "9. Admin → System: Generate Invoice() → Display Invoice()",
        "10. Admin → System: Logout()",
    ]
    s = add_multiline_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.2),
                      seq_lines, font_size=16, color=LIGHT_GRAY)
    add_entrance_animation(slide, s, 300)


# ═══════════════════════════════════════════════════════════════════
#  SLIDE 19 — THANK YOU
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.8), W, Inches(2))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

s1 = add_text_box(slide, Inches(1), Inches(3), Inches(11), Inches(1.2),
             "THANK YOU", font_size=60, bold=True, color=WHITE,
             alignment=PP_ALIGN.CENTER)
s2 = add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.8),
             "Nathkrupa Tours & Travels — Project Documentation",
             font_size=20, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
add_entrance_animation(slide, s1, 0)
add_entrance_animation(slide, s2, 500)
add_slide_number(slide, 19)


# ═══════════════════════════════════════════════════════════════════
#  ADD SLIDE TRANSITIONS (Fade) to all slides
# ═══════════════════════════════════════════════════════════════════
import lxml.etree as etree

for slide in prs.slides:
    transition = etree.SubElement(
        slide._element,
        '{http://schemas.openxmlformats.org/presentationml/2006/main}transition',
        attrib={'spd': 'med', 'advClick': '1'}
    )
    etree.SubElement(transition, '{http://schemas.openxmlformats.org/presentationml/2006/main}fade',
                     attrib={'thruBlk': '1'})


# ═══════════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════════
prs.save(SAVE_PATH)
print(f"\n[OK] Presentation saved successfully!")
print(f"Location: {SAVE_PATH}")
print(f"Total slides: {len(prs.slides)}")
